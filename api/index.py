import datetime
import hmac
import math
import os
import sys
import zipfile
import uuid
import tarfile
import tempfile
import requests
from functools import wraps
from flask import Flask, Response, render_template, request, send_file, jsonify, session, redirect, url_for, g
from io import BytesIO
from PyPDF2 import PdfMerger, PdfReader, PdfWriter
from pdf2docx import Converter as PdfToDocxConverter
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from reportlab.pdfgen import canvas as pdf_canvas
from reportlab.lib.utils import ImageReader
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
import reportlab as _reportlab

# Helvetica (PDF base14) WinAnsiEncoding kullanır ve Türkçe İ/ı/ğ/Ş gibi
# karakterleri içermez, bu yüzden kutu (.notdef) olarak basılır.
# reportlab'ın kendi paketiyle gelen Bitstream Vera Sans tam Unicode
# glyph setine sahip olduğundan Türkçe karakterleri doğru render eder.
PDF_TEXT_FONT = "PPTXUnicode"
pdfmetrics.registerFont(
    TTFont(PDF_TEXT_FONT, os.path.join(os.path.dirname(_reportlab.__file__), "fonts", "Vera.ttf"))
)

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from i18n import SUPPORTED_LANGS, init_i18n, t
import analytics
from travel_phrases import EMERGENCY_CATEGORIES, PHRASEBOOK_CATEGORIES
from travel_places import NEARBY_CATEGORIES
from travel_currencies import CURRENCIES, QUICK_REFERENCE_CODES

# Overpass (OpenStreetMap) etiket filtrelerine göre desteklenen kategoriler.
# Yalnızca burada olan kategoriler için "Yakında Ara" gerçek bir liste gösterir;
# "consulate" gibi OSM'de tutarlı/ülkeye özel etiketlenmeyen kategoriler
# doğrudan Google Haritalar aramasına yönlendirmeye devam eder.
OVERPASS_CATEGORY_FILTERS = {
    "hospital": ['["amenity"="hospital"]'],
    "pharmacy": ['["amenity"="pharmacy"]'],
    "police": ['["amenity"="police"]'],
    "restaurant": ['["amenity"="restaurant"]'],
    "atm": ['["amenity"="atm"]'],
    "gas": ['["amenity"="fuel"]'],
    "transit": [
        '["highway"="bus_stop"]',
        '["amenity"="bus_station"]',
        '["railway"="station"]',
        '["railway"="tram_stop"]',
    ],
}

# category -> çeviri anahtarı (isimsiz OSM sonuçlarında yedek isim göstermek için)
_CATEGORY_LABEL_KEYS = {key: label_key for key, _icon, label_key, _term in NEARBY_CATEGORIES}


def _haversine_m(lat1, lng1, lat2, lng2):
    r = 6371000
    phi1, phi2 = math.radians(lat1), math.radians(lat2)
    d_phi = math.radians(lat2 - lat1)
    d_lambda = math.radians(lng2 - lng1)
    a = math.sin(d_phi / 2) ** 2 + math.cos(phi1) * math.cos(phi2) * math.sin(d_lambda / 2) ** 2
    return 2 * r * math.asin(math.sqrt(a))

app = Flask(__name__)
app.secret_key = os.environ.get("FLASK_SECRET_KEY", "dev-insecure-change-me-in-prod")
init_i18n(app)

ADMIN_PASSWORD = os.environ.get("ADMIN_PASSWORD")

# Vercel prod/preview/dev ortamlarında otomatik olarak set edilir; yerel
# `flask run` sırasında tanımlı değildir. Cookie'lerin Secure bayrağını
# yerel http testini bozmadan yalnızca gerçek (https) ortamda açmak için kullanılır.
IS_VERCEL = bool(os.environ.get("VERCEL_ENV"))
app.config["SESSION_COOKIE_SECURE"] = IS_VERCEL

# Sitemap için site içindeki tüm statik/GET sayfa yolları
SITEMAP_ROUTES = [
    ("/", 1.0),
    ("/araclar", 0.95),
    ("/mikrotik", 0.8),
    ("/mikrotik/nasil-kullanilir", 0.5),
    ("/mikrotik/subnet", 0.6),
    ("/mikrotik/ip-guide", 0.5),
    ("/mikrotik/dhcp-server", 0.5),
    ("/mikrotik/interface-ip", 0.5),
    ("/mikrotik/port-forward", 0.5),
    ("/mikrotik/src-nat", 0.5),
    ("/mikrotik/routing", 0.5),
    ("/mikrotik/load-balance", 0.5),
    ("/mikrotik/firewall", 0.5),
    ("/mikrotik/dns-security", 0.5),
    ("/mikrotik/port-services", 0.5),
    ("/mikrotik/vpn", 0.5),
    ("/mikrotik/netwatch", 0.5),
    ("/mikrotik/log-helper", 0.5),
    ("/mikrotik/bandwidth-pcq", 0.5),
    ("/mikrotik/fasttrack", 0.5),
    ("/mikrotik/wifi-channels", 0.5),
    ("/mikrotik/hardware-selector", 0.5),
    ("/mikrotik/hotspot", 0.5),
    ("/video-tools", 0.7),
    ("/audio-tools", 0.7),
    ("/ai-tools", 0.7),
    ("/guncelleme-notlari", 0.3),
    ("/pdf-tools", 0.8),
    ("/pdf-tools/split", 0.6),
    ("/pdf-tools/protect", 0.6),
    ("/pdf-tools/unlock", 0.6),
    ("/pdf-tools/rotate", 0.6),
    ("/pdf-tools/page-numbers", 0.6),
    ("/pdf-tools/compress", 0.6),
    ("/pdf-tools/pdf-to-jpg", 0.6),
    ("/pdf-tools/pdf-to-png", 0.6),
    ("/pdf-tools/jpg-to-pdf", 0.6),
    ("/pdf-tools/png-to-pdf", 0.6),
    ("/pdf-tools/pdf-to-html", 0.6),
    ("/pdf-tools/word-to-pdf", 0.6),
    ("/pdf-tools/pdf-to-word", 0.6),
    ("/pdf-tools/excel-to-pdf", 0.6),
    ("/pdf-tools/pdf-to-excel", 0.6),
    ("/pdf-tools/ppt-to-pdf", 0.6),
    ("/audio-tools/trim", 0.5),
    ("/audio-tools/record", 0.5),
    ("/audio-tools/volume", 0.5),
    ("/audio-tools/speed", 0.5),
    ("/audio-tools/pitch", 0.5),
    ("/audio-tools/equalizer", 0.5),
    ("/audio-tools/joiner", 0.5),
    ("/audio-tools/reverse", 0.5),
    ("/video-tools/screen-record", 0.5),
    ("/video-tools/rotate", 0.5),
    ("/video-tools/trim", 0.5),
    ("/video-tools/merge", 0.5),
    ("/video-tools/editor", 0.5),
    ("/video-tools/crop", 0.5),
    ("/video-tools/loop", 0.5),
    ("/video-tools/video-volume", 0.5),
    ("/video-tools/video-speed", 0.5),
    ("/video-tools/add-text", 0.5),
    ("/video-tools/add-image", 0.5),
    ("/video-tools/stabilize", 0.5),
    ("/video-tools/remove-logo", 0.5),
    ("/video-tools/add-audio", 0.5),
    ("/video-tools/resize", 0.5),
    ("/video-tools/text-to-speech", 0.5),
    ("/video-tools/record-camera", 0.5),
    ("/video-tools/gif", 0.5),
    ("/qr-kod", 0.6),
    ("/convert/image", 0.6),
    ("/convert/image-compress", 0.6),
    ("/convert/audio", 0.6),
    ("/convert/video", 0.6),
    ("/convert/extract", 0.6),
    ("/convert/document", 0.6),
    ("/convert/ebook", 0.6),
    ("/convert/font", 0.6),
    ("/convert/archive", 0.6),
    ("/seyahat", 0.7),
    ("/seyahat/acil-ceviri", 0.7),
    ("/seyahat/yakin-yerler", 0.7),
    ("/seyahat/konusma-kilavuzu", 0.7),
    ("/seyahat/doviz", 0.7),
]

_TRACKED_PATHS = {path for path, _ in SITEMAP_ROUTES}
_VISITOR_COOKIE = "vid"

# Ana sayfadaki "En Çok Kullanılan Araçlar" bölümü için path -> kart meta verisi.
# Yalnızca _TRACKED_PATHS içindeki (ziyaret sayısı tutulan) yollar buradadır.
TOOL_CARDS = {
    "/ai-tools": ("ai", "🤖", "active", "home.ai.compare.title", "home.ai.compare.desc"),
    "/pdf-tools": ("pdf", "📚", "active", "home.pdf.merge.title", "home.pdf.merge.desc"),
    "/pdf-tools/split": ("pdf", "✂️", "active", "home.pdf.split.title", "home.pdf.split.desc"),
    "/pdf-tools/protect": ("pdf", "🔒", "active", "home.pdf.protect.title", "home.pdf.protect.desc"),
    "/pdf-tools/unlock": ("pdf", "🔓", "active", "home.pdf.unlock.title", "home.pdf.unlock.desc"),
    "/pdf-tools/rotate": ("pdf", "🔄", "active", "home.pdf.rotate.title", "home.pdf.rotate.desc"),
    "/pdf-tools/page-numbers": ("pdf", "🔢", "active", "home.pdf.page_numbers.title", "home.pdf.page_numbers.desc"),
    "/pdf-tools/compress": ("pdf", "🗜️", "active", "home.pdf.compress.title", "home.pdf.compress.desc"),
    "/pdf-tools/pdf-to-jpg": ("pdf", "🖼️", "active", "home.pdf.to_jpg.title", "home.pdf.to_jpg.desc"),
    "/pdf-tools/pdf-to-png": ("pdf", "🖼️", "active", "home.pdf.to_png.title", "home.pdf.to_png.desc"),
    "/pdf-tools/jpg-to-pdf": ("pdf", "📄", "active", "home.pdf.jpg_to_pdf.title", "home.pdf.jpg_to_pdf.desc"),
    "/pdf-tools/png-to-pdf": ("pdf", "📄", "active", "home.pdf.png_to_pdf.title", "home.pdf.png_to_pdf.desc"),
    "/pdf-tools/pdf-to-html": ("pdf", "🌐", "active", "home.pdf.to_html.title", "home.pdf.to_html.desc"),
    "/pdf-tools/word-to-pdf": ("pdf", "📝", "active", "home.pdf.word_to_pdf.title", "home.pdf.word_to_pdf.desc"),
    "/pdf-tools/pdf-to-word": ("pdf", "📝", "active", "home.pdf.to_word.title", "home.pdf.to_word.desc"),
    "/pdf-tools/excel-to-pdf": ("pdf", "📊", "active", "home.pdf.excel_to_pdf.title", "home.pdf.excel_to_pdf.desc"),
    "/pdf-tools/pdf-to-excel": ("pdf", "📊", "active", "home.pdf.to_excel.title", "home.pdf.to_excel.desc"),
    "/pdf-tools/ppt-to-pdf": ("pdf", "📉", "active", "home.pdf.ppt_to_pdf.title", "home.pdf.ppt_to_pdf.desc"),
    "/convert/image": ("converter", "🖼️", "active", "home.conv.image.title", "home.conv.image.desc"),
    "/convert/image-compress": ("converter", "🗜️", "active", "home.conv.image_compress.title", "home.conv.image_compress.desc"),
    "/convert/audio": ("converter", "🎵", "active", "home.conv.audio.title", "home.conv.audio.desc"),
    "/convert/video": ("converter", "🎬", "active", "home.conv.video.title", "home.conv.video.desc"),
    "/convert/extract": ("converter", "🗄️", "active", "home.conv.extract.title", "home.conv.extract.desc"),
    "/convert/document": ("converter", "📝", "active", "home.conv.document.title", "home.conv.document.desc"),
    "/convert/ebook": ("converter", "📚", "active", "home.conv.ebook.title", "home.conv.ebook.desc"),
    "/convert/font": ("converter", "🔤", "active", "home.conv.font.title", "home.conv.font.desc"),
    "/convert/archive": ("converter", "🗄️", "active", "home.conv.archive.title", "home.conv.archive.desc"),
    "/video-tools/screen-record": ("video", "🎥", "beta", "home.video.screen_record.title", "home.video.screen_record.desc"),
    "/video-tools/trim": ("video", "✂️", "beta", "home.video.trim.title", "home.video.trim.desc"),
    "/video-tools/merge": ("video", "🔀", "beta", "home.video.merge.title", "home.video.merge.desc"),
    "/video-tools/rotate": ("video", "🔄", "beta", "home.video.rotate.title", "home.video.rotate.desc"),
    "/video-tools/editor": ("video", "🎬", "beta", "home.video.editor.title", "home.video.editor.desc"),
    "/video-tools/add-image": ("video", "🖼️", "beta", "home.video.add_image.title", "home.video.add_image.desc"),
    "/video-tools/resize": ("video", "📐", "beta", "home.video.resize.title", "home.video.resize.desc"),
    "/video-tools/add-text": ("video", "✍️", "beta", "home.video.add_text.title", "home.video.add_text.desc"),
    "/video-tools/loop": ("video", "🔁", "beta", "home.video.loop.title", "home.video.loop.desc"),
    "/video-tools/text-to-speech": ("video", "🗣️", "beta", "home.video.tts.title", "home.video.tts.desc"),
    "/video-tools/remove-logo": ("video", "🚫", "beta", "home.video.remove_logo.title", "home.video.remove_logo.desc"),
    "/video-tools/video-volume": ("video", "🔊", "beta", "home.video.volume.title", "home.video.volume.desc"),
    "/video-tools/crop": ("video", "✂️", "beta", "home.video.crop.title", "home.video.crop.desc"),
    "/video-tools/video-speed": ("video", "⚡", "beta", "home.video.speed.title", "home.video.speed.desc"),
    "/video-tools/stabilize": ("video", "⚖️", "beta", "home.video.stabilize.title", "home.video.stabilize.desc"),
    "/video-tools/add-audio": ("video", "🎵", "beta", "home.video.add_audio.title", "home.video.add_audio.desc"),
    "/video-tools/record-camera": ("video", "📹", "beta", "home.video.record_camera.title", "home.video.record_camera.desc"),
    "/video-tools/gif": ("video", "🎞️", "beta", "home.video.gif.title", "home.video.gif.desc"),
    "/audio-tools/trim": ("audio", "✂️", "beta", "home.audio.trim.title", "home.audio.trim.desc"),
    "/audio-tools/record": ("audio", "🎙️", "beta", "home.audio.record.title", "home.audio.record.desc"),
    "/audio-tools/volume": ("audio", "📢", "beta", "home.audio.volume.title", "home.audio.volume.desc"),
    "/audio-tools/speed": ("audio", "⚡", "beta", "home.audio.speed.title", "home.audio.speed.desc"),
    "/audio-tools/pitch": ("audio", "🎼", "beta", "home.audio.pitch.title", "home.audio.pitch.desc"),
    "/audio-tools/equalizer": ("audio", "🎚️", "beta", "home.audio.equalizer.title", "home.audio.equalizer.desc"),
    "/audio-tools/joiner": ("audio", "🔗", "beta", "home.audio.joiner.title", "home.audio.joiner.desc"),
    "/audio-tools/reverse": ("audio", "🔄", "beta", "home.audio.reverse.title", "home.audio.reverse.desc"),
    "/qr-kod": ("utility", "📱", "active", "home.utility.qr.title", "home.utility.qr.desc"),
}
_POPULAR_TOOLS_LIMIT = 6


def _popular_tools():
    """En çok ziyaret edilen ve TOOL_CARDS'ta karşılığı olan araçları döner.
    KV yapılandırılmamışsa (yerelde) boş liste döner, bölüm gizlenir."""
    stats = analytics.get_stats()
    if not stats.get("available"):
        return []
    popular = []
    for path, _count in stats["routes"]:
        card = TOOL_CARDS.get(path)
        if not card:
            continue
        cat, icon, badge, title_key, desc_key = card
        popular.append({
            "path": path,
            "cat": cat,
            "icon": icon,
            "badge": badge,
            "title_key": title_key,
            "desc_key": desc_key,
        })
        if len(popular) >= _POPULAR_TOOLS_LIMIT:
            break
    return popular


def _client_ip():
    forwarded = request.headers.get("X-Forwarded-For", "")
    if forwarded:
        return forwarded.split(",")[0].strip()
    return request.remote_addr or "unknown"


def _log_action(action):
    """Bir aracın gerçekten kullanıldığını kaydeder (dosya/belge içeriği değil,
    yalnızca kim/ne zaman/hangi işlem)."""
    analytics.log_event(_client_ip(), request.path, "action", action=action)


ACTION_LABELS = {
    "pdf_merge": "PDF Birleştirme",
    "pdf_split": "PDF Bölme",
    "pdf_protect": "PDF Şifreleme",
    "pdf_unlock": "PDF Kilit Açma",
    "zip_extract": "ZIP Çıkartma",
    "pdf_to_word": "PDF -> Word",
    "ppt_to_pdf": "PPT -> PDF",
}


EMU_PER_POINT = 12700


def _emu_to_pt(emu):
    return (emu or 0) / EMU_PER_POINT


def _autofit_font_scale(text_frame):
    """PowerPoint'in 'Sığdırmak için metni küçült' (normAutofit) ayarını okur.
    PowerPoint bu ayar varsa yazıyı ekranda gösterirken küçültür; biz bunu
    uygulamazsak metin olduğundan büyük çizilip şekil sınırlarını taşar."""
    try:
        body_pr = text_frame._txBody.bodyPr
        norm_autofit = body_pr.find(qn("a:normAutofit"))
        if norm_autofit is not None:
            scale_attr = norm_autofit.get("fontScale")
            if scale_attr:
                return int(scale_attr) / 100000.0
    except Exception:
        pass
    return 1.0


def _shrink_scale_to_fit(paragraph_sizes, font_name, max_width, max_height, min_scale=0.35):
    """PowerPoint'in normAutofit'i XML'e yazmadığı (ör. programatik/dönüştürülmüş
    dosyalar) durumlarda bile metnin şekil sınırlarını taşmasını önlemek için
    kendi 'sığdırmak üzere küçült' hesabımızı yapar: toplam satır yüksekliği
    şeklin yüksekliğini aşıyorsa font oranı kademeli olarak küçültülür."""
    if max_height <= 0:
        return 1.0
    scale = 1.0
    while scale > min_scale:
        total_height = 0.0
        for text, base_size in paragraph_sizes:
            size = base_size * scale
            if text.strip():
                total_height += len(_wrap_text_lines(text, font_name, size, max_width)) * size * 1.2
            else:
                total_height += size * 1.2
        if total_height <= max_height:
            return scale
        scale -= 0.05
    return min_scale


def _wrap_text_lines(text, font_name, font_size, max_width):
    """Metni verilen genişliğe sığacak şekilde kelime kelime satırlara böler.
    reportlab drawString satır kaydırma yapmadığından, şekil genişliğini
    aşan metinler sayfa dışına taşıp görünmez oluyordu."""
    if max_width <= 0:
        return [text]
    words = text.split(" ")
    lines = []
    current = ""
    for word in words:
        candidate = f"{current} {word}".strip() if current else word
        if not current or pdfmetrics.stringWidth(candidate, font_name, font_size) <= max_width:
            current = candidate
        else:
            lines.append(current)
            current = word
    if current:
        lines.append(current)
    return lines


def _convert_pptx_to_pdf(pptx_stream):
    """python-pptx ile slaytların gerçek içeriğini (metin, görsel, tablo)
    okuyup reportlab ile aynı konum/boyutta bir PDF sayfasına çizer.
    Grup şekillerin içi ve gelişmiş biçimlendirme (renk, hizalama) kapsam dışıdır."""
    prs = Presentation(pptx_stream)
    slide_w = _emu_to_pt(prs.slide_width)
    slide_h = _emu_to_pt(prs.slide_height)

    output = BytesIO()
    c = pdf_canvas.Canvas(output, pagesize=(slide_w, slide_h))

    for slide in prs.slides:
        for shape in slide.shapes:
            try:
                left = _emu_to_pt(shape.left)
                top = _emu_to_pt(shape.top)
                width = _emu_to_pt(shape.width)
                height = _emu_to_pt(shape.height)
                top_y = slide_h - top

                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    img_stream = BytesIO(shape.image.blob)
                    c.drawImage(
                        ImageReader(img_stream),
                        left, top_y - height, width=width, height=height,
                        preserveAspectRatio=True, mask="auto",
                    )
                elif shape.has_table:
                    table = shape.table
                    rows, cols = len(table.rows), len(table.columns)
                    cell_h = height / max(rows, 1)
                    cell_w = width / max(cols, 1)
                    for r in range(rows):
                        for col in range(cols):
                            cx = left + col * cell_w
                            cy = top_y - (r + 1) * cell_h
                            c.setFont(PDF_TEXT_FONT, 9)
                            c.drawString(cx + 2, cy + 2, table.cell(r, col).text[:80])
                elif shape.has_text_frame:
                    max_width = width if width > 0 else (slide_w - left)
                    font_scale = _autofit_font_scale(shape.text_frame)

                    paragraph_sizes = []
                    for para in shape.text_frame.paragraphs:
                        text = "".join(run.text for run in para.runs) or para.text
                        font_size = 14
                        if para.runs and para.runs[0].font.size:
                            font_size = para.runs[0].font.size.pt
                        paragraph_sizes.append((text, font_size * font_scale))

                    # normAutofit XML'de yoksa (ör. programatik olarak üretilmiş
                    # dosyalarda) kendi hesabımızla şekil yüksekliğine sığdırıyoruz.
                    fit_scale = _shrink_scale_to_fit(paragraph_sizes, PDF_TEXT_FONT, max_width, height)

                    line_y = top_y - 14
                    for text, base_size in paragraph_sizes:
                        font_size = base_size * fit_scale
                        if text.strip():
                            c.setFont(PDF_TEXT_FONT, font_size)
                            for line in _wrap_text_lines(text, PDF_TEXT_FONT, font_size, max_width):
                                c.drawString(left, line_y, line)
                                line_y -= font_size * 1.2
                        else:
                            line_y -= font_size * 1.2
            except Exception:
                continue
        c.showPage()

    c.save()
    output.seek(0)
    return output


def _action_label(action):
    if action in ACTION_LABELS:
        return ACTION_LABELS[action]
    if action.startswith("archive_convert_to_"):
        return f"Arşiv Dönüştürme (→ {action[len('archive_convert_to_'):].upper()})"
    return action


@app.before_request
def _track_pageview():
    if request.method != "GET" or request.path not in _TRACKED_PATHS:
        return
    vid = request.cookies.get(_VISITOR_COOKIE)
    if not vid:
        vid = uuid.uuid4().hex
        g.new_vid = vid
    analytics.track_pageview(request.path, vid, _client_ip())


@app.after_request
def _persist_visitor_cookie(response):
    vid = getattr(g, "new_vid", None)
    if vid:
        response.set_cookie(
            _VISITOR_COOKIE,
            vid,
            max_age=60 * 60 * 24 * 365 * 2,
            samesite="Lax",
            httponly=True,
            secure=IS_VERCEL,
        )
    return response


@app.route("/sitemap.xml")
def sitemap():
    root = request.host_url.rstrip("/")
    entries = []
    for path, priority in SITEMAP_ROUTES:
        url = f"{root}{path}"
        alt_links = "".join(
            f'<xhtml:link rel="alternate" hreflang="{lang}" href="{root}{path}?lang={lang}" />'
            for lang in SUPPORTED_LANGS
        )
        entries.append(
            f"<url><loc>{url}</loc>{alt_links}<priority>{priority}</priority></url>"
        )
    xml = (
        '<?xml version="1.0" encoding="UTF-8"?>'
        '<urlset xmlns="http://www.sitemaps.org/schemas/sitemap/0.9" '
        'xmlns:xhtml="http://www.w3.org/1999/xhtml">'
        + "".join(entries)
        + "</urlset>"
    )
    return Response(xml, mimetype="application/xml")


@app.route("/robots.txt")
def robots_txt():
    root = request.host_url.rstrip("/")
    lines = [
        "User-agent: *",
        "Disallow: /convert/extract/download",
        "Disallow: /admin",
        f"Sitemap: {root}/sitemap.xml",
    ]
    return Response("\n".join(lines), mimetype="text/plain")


def _admin_required(view):
    @wraps(view)
    def wrapped(*args, **kwargs):
        if not session.get("is_admin"):
            return redirect(url_for("admin_login"))
        return view(*args, **kwargs)

    return wrapped


@app.route("/admin/login", methods=["GET", "POST"])
def admin_login():
    if session.get("is_admin"):
        return redirect(url_for("admin_dashboard"))

    error = None
    if request.method == "POST":
        client_id = _client_ip()
        lock_seconds = analytics.check_login_lock(client_id)
        if lock_seconds > 0:
            error = f"Çok fazla hatalı deneme. {lock_seconds} saniye sonra tekrar dene."
        else:
            password = request.form.get("password", "")
            if ADMIN_PASSWORD and hmac.compare_digest(password, ADMIN_PASSWORD):
                analytics.clear_login_attempts(client_id)
                session["is_admin"] = True
                return redirect(url_for("admin_dashboard"))
            analytics.record_failed_login(client_id)
            error = "Şifre yanlış." if ADMIN_PASSWORD else "ADMIN_PASSWORD ortam değişkeni tanımlı değil."
    return render_template("admin_login.html", error=error)


@app.route("/admin/logout")
def admin_logout():
    session.pop("is_admin", None)
    return redirect(url_for("admin_login"))


@app.route("/admin")
@_admin_required
def admin_dashboard():
    stats = analytics.get_stats()
    stats["actions"] = [
        {"label": _action_label(action), "raw": action, "count": count}
        for action, count in stats.get("actions", [])
    ]
    location_stats = analytics.get_location_stats()
    return render_template("admin_dashboard.html", stats=stats, location_stats=location_stats)


@app.route("/admin/logs")
@_admin_required
def admin_logs():
    try:
        day = datetime.datetime.strptime(request.args.get("day", ""), "%Y-%m-%d").date()
    except ValueError:
        day = datetime.datetime.utcnow().date()

    day_str = day.strftime("%Y-%m-%d")
    events = analytics.get_log(day_str)
    return render_template(
        "admin_logs.html",
        day=day_str,
        prev_day=(day - datetime.timedelta(days=1)).strftime("%Y-%m-%d"),
        next_day=(day + datetime.timedelta(days=1)).strftime("%Y-%m-%d"),
        events=events,
        analytics_available=analytics.is_configured(),
    )


# ZIP bomb koruması sabitleri
_ZIP_MAX_FILES = 30
_ZIP_MAX_TOTAL_BYTES = 100 * 1024 * 1024   # 100 MB açılmış toplam boyut
_ZIP_MAX_FILE_BYTES = 50 * 1024 * 1024    # 50 MB tek dosya
_ZIP_MAX_RATIO = 100                       # sıkıştırma oranı üst sınırı


def _check_zip_safety(zip_ref: zipfile.ZipFile):
    """Zip bomb kontrolü. Sorun varsa (False, hata_mesajı) döner."""
    members = [i for i in zip_ref.infolist() if not i.is_dir()]

    if len(members) > _ZIP_MAX_FILES:
        return False, t(
            "err.zip_too_many_files",
            max_files=_ZIP_MAX_FILES,
            count=len(members),
        )

    total = 0
    for info in members:
        if info.file_size > _ZIP_MAX_FILE_BYTES:
            return False, t(
                "err.zip_file_too_large",
                filename=info.filename,
                max_mb=_ZIP_MAX_FILE_BYTES // (1024 * 1024),
            )
        if info.compress_size > 0 and (info.file_size / info.compress_size) > _ZIP_MAX_RATIO:
            return False, t(
                "err.zip_suspicious_ratio",
                ratio=info.file_size // max(info.compress_size, 1),
            )
        total += info.file_size
        if total > _ZIP_MAX_TOTAL_BYTES:
            return False, t(
                "err.zip_total_too_large",
                max_mb=_ZIP_MAX_TOTAL_BYTES // (1024 * 1024),
            )

    return True, None


@app.route("/")
def landing():
    return render_template("landing.html")


@app.route("/araclar")
def tools_home():
    return render_template("index.html", popular_tools=_popular_tools())


@app.route("/mikrotik")
def mikrotik():
    return render_template("mikrotik.html", active_page="dashboard")


@app.route("/mikrotik/nasil-kullanilir")
def mikrotik_guide():
    return render_template("mikrotik_guide.html", active_page="nasil-kullanilir")


@app.route("/mikrotik/subnet")
def mikrotik_subnet():
    return render_template("mikrotik_subnet.html", active_page="subnet")


@app.route("/mikrotik/ip-guide")
def mikrotik_ip_guide():
    return render_template("mikrotik_ip_guide.html", active_page="ip-guide")


@app.route("/mikrotik/dhcp-server")
def mikrotik_dhcp():
    return render_template("mikrotik_dhcp.html", active_page="dhcp-server")


@app.route("/mikrotik/interface-ip")
def mikrotik_interface():
    return render_template("mikrotik_interface.html", active_page="interface-ip")


@app.route("/mikrotik/port-forward")
def mikrotik_port_forward():
    return render_template("mikrotik_port_forward.html", active_page="port-forward")


@app.route("/mikrotik/src-nat")
def mikrotik_src_nat():
    return render_template("mikrotik_src_nat.html", active_page="src-nat")


@app.route("/mikrotik/routing")
def mikrotik_routing():
    return render_template("mikrotik_routing.html", active_page="routing")


@app.route("/mikrotik/load-balance")
def mikrotik_load_balance():
    return render_template("mikrotik_load_balance.html", active_page="load-balance")


@app.route("/mikrotik/firewall")
def mikrotik_firewall():
    return render_template("mikrotik_firewall.html", active_page="firewall")


@app.route("/mikrotik/dns-security")
def mikrotik_dns():
    return render_template("mikrotik_dns.html", active_page="dns-security")


@app.route("/mikrotik/port-services")
def mikrotik_port_services():
    return render_template("mikrotik_port_services.html", active_page="port-services")


@app.route("/mikrotik/vpn")
def mikrotik_vpn():
    return render_template("mikrotik_vpn.html", active_page="vpn")


@app.route("/mikrotik/netwatch")
def mikrotik_netwatch():
    return render_template("mikrotik_netwatch.html", active_page="netwatch")


@app.route("/mikrotik/log-helper")
def mikrotik_log():
    return render_template("mikrotik_log.html", active_page="log-helper")


@app.route("/mikrotik/bandwidth-pcq")
def mikrotik_bandwidth():
    return render_template("mikrotik_bandwidth.html", active_page="bandwidth-pcq")


@app.route("/mikrotik/fasttrack")
def mikrotik_fasttrack():
    return render_template("mikrotik_fasttrack.html", active_page="fasttrack")


@app.route("/mikrotik/wifi-channels")
def mikrotik_wifi():
    return render_template("mikrotik_wifi.html", active_page="wifi-channels")


@app.route("/mikrotik/hardware-selector")
def mikrotik_hardware():
    return render_template("mikrotik_hardware.html", active_page="hardware-selector")


@app.route("/mikrotik/hotspot")
def mikrotik_hotspot():
    return render_template("mikrotik_hotspot.html", active_page="hotspot")


@app.route("/seyahat")
def travel_dashboard():
    return render_template("travel_dashboard.html", active_page="dashboard")


@app.route("/seyahat/acil-ceviri")
def travel_emergency_translate():
    return render_template(
        "travel_emergency.html",
        active_page="acil-ceviri",
        categories=EMERGENCY_CATEGORIES,
    )


@app.route("/seyahat/yakin-yerler")
def travel_nearby():
    return render_template(
        "travel_nearby.html",
        active_page="yakin-yerler",
        categories=NEARBY_CATEGORIES,
    )


@app.route("/seyahat/konusma-kilavuzu")
def travel_phrasebook():
    return render_template(
        "travel_phrasebook.html",
        active_page="konusma-kilavuzu",
        categories=PHRASEBOOK_CATEGORIES,
    )


@app.route("/seyahat/doviz")
def travel_currency():
    return render_template(
        "travel_currency.html",
        active_page="doviz",
        currencies=CURRENCIES,
        quick_codes=QUICK_REFERENCE_CODES,
    )


@app.route("/seyahat/reverse-geocode")
def travel_reverse_geocode():
    """Koordinatı okunabilir adrese çevirir (OpenStreetMap Nominatim).
    Konum bu sunucu üzerinden tek seferlik geçer, hiçbir yerde saklanmaz."""
    try:
        lat = float(request.args.get("lat", ""))
        lng = float(request.args.get("lng", ""))
    except (TypeError, ValueError):
        return jsonify({"error": "invalid_coords"}), 400

    if not (-90 <= lat <= 90) or not (-180 <= lng <= 180):
        return jsonify({"error": "invalid_coords"}), 400

    if not analytics.try_acquire_ip_slot(_client_ip(), "reverse-geocode", 15):
        return jsonify({"error": "rate_limited"}), 429

    if not analytics.try_acquire_geocode_slot():
        return jsonify({"error": "rate_limited"}), 429

    analytics.log_location_query("reverse_geocode")

    try:
        resp = requests.get(
            "https://nominatim.openstreetmap.org/reverse",
            params={"lat": lat, "lon": lng, "format": "jsonv2", "zoom": 18, "addressdetails": 0},
            headers={"User-Agent": "toolboxquick-travel-assistant/1.0 (contact: kaanbyzt07@gmail.com)"},
            timeout=4,
        )
        resp.raise_for_status()
        data = resp.json()
    except (requests.RequestException, ValueError):
        return jsonify({"error": "geocode_failed"}), 502

    address = data.get("display_name") if isinstance(data, dict) else None
    if not address:
        return jsonify({"error": "not_found"}), 404

    return jsonify({"address": address})


@app.route("/seyahat/nearby-places")
def travel_nearby_places():
    """Bir kategori için yakındaki gerçek yerleri (OpenStreetMap/Overpass API)
    isim, mesafe ve adres bilgisiyle listeler. Konum bu sunucu üzerinden tek
    seferlik geçer, hiçbir yerde saklanmaz."""
    category = request.args.get("category", "")
    filters = OVERPASS_CATEGORY_FILTERS.get(category)
    if not filters:
        return jsonify({"error": "unsupported_category"}), 400

    try:
        lat = float(request.args.get("lat", ""))
        lng = float(request.args.get("lng", ""))
    except (TypeError, ValueError):
        return jsonify({"error": "invalid_coords"}), 400

    if not (-90 <= lat <= 90) or not (-180 <= lng <= 180):
        return jsonify({"error": "invalid_coords"}), 400

    if not analytics.try_acquire_ip_slot(_client_ip(), "nearby-places", 15):
        return jsonify({"error": "rate_limited"}), 429

    if not analytics.try_acquire_overpass_slot():
        return jsonify({"error": "rate_limited"}), 429

    analytics.log_location_query("nearby_places")

    radius_m = 6000
    clauses = "".join(
        f'node{f}(around:{radius_m},{lat},{lng});way{f}(around:{radius_m},{lat},{lng});'
        for f in filters
    )
    query = f"[out:json][timeout:15];({clauses});out center 30;"

    try:
        resp = requests.post(
            "https://overpass-api.de/api/interpreter",
            data={"data": query},
            headers={"User-Agent": "toolboxquick-travel-assistant/1.0 (contact: kaanbyzt07@gmail.com)"},
            timeout=16,
        )
        resp.raise_for_status()
        data = resp.json()
    except (requests.RequestException, ValueError):
        return jsonify({"error": "lookup_failed"}), 502

    fallback_name = t(_CATEGORY_LABEL_KEYS.get(category, ""))
    results = []
    for el in data.get("elements", []) if isinstance(data, dict) else []:
        tags = el.get("tags") or {}
        name = tags.get("name") or tags.get("operator") or fallback_name
        if not name:
            continue

        if el.get("type") == "node":
            p_lat, p_lng = el.get("lat"), el.get("lon")
        else:
            center = el.get("center") or {}
            p_lat, p_lng = center.get("lat"), center.get("lon")

        if not isinstance(p_lat, (int, float)) or not isinstance(p_lng, (int, float)):
            continue

        street = " ".join(p for p in [tags.get("addr:street"), tags.get("addr:housenumber")] if p)
        city = tags.get("addr:city") or tags.get("addr:suburb") or tags.get("addr:district")
        address = ", ".join(p for p in [street, city] if p) or None

        results.append({
            "name": name,
            "lat": p_lat,
            "lng": p_lng,
            "distance_m": round(_haversine_m(lat, lng, p_lat, p_lng)),
            "address": address,
            "phone": tags.get("phone") or tags.get("contact:phone"),
        })

    results.sort(key=lambda r: r["distance_m"])
    return jsonify({"results": results[:15]})


@app.route("/video-tools")
def video_tools():
    return render_template("video_tools.html")


@app.route("/audio-tools")
def audio_tools():
    return render_template("audio_tools.html")


@app.route("/ai-tools")
def ai_tools():
    return render_template("ai_tools.html")


@app.route("/guncelleme-notlari")
def changelog():
    return render_template("changelog.html")


@app.route("/pdf-tools", methods=["GET", "POST"])
def pdf_tools():
    if request.method == "GET":
        return render_template("pdf_tools.html", download_url=None, error=None)

    files = request.files.getlist("pdf_files")
    pdf_files = [f for f in files if f and f.filename.lower().endswith(".pdf")]

    if len(pdf_files) < 2:
        return render_template(
            "pdf_tools.html",
            download_url=None,
            error=t("err.min_two_pdfs"),
        )

    merger = PdfMerger()
    for f in pdf_files:
        merger.append(f.stream)

    output = BytesIO()
    merger.write(output)
    merger.close()
    output.seek(0)

    _log_action("pdf_merge")
    return send_file(
        output,
        mimetype="application/pdf",
        as_attachment=True,
        download_name="toolboxquick-merged.pdf",
    )


@app.route("/pdf-tools/split", methods=["GET", "POST"])
def pdf_split():
    if request.method == "GET":
        return render_template("pdf_split.html", error=None)

    file = request.files.get("pdf_file")
    range_str = request.form.get("page_range", "").strip()

    if not file or not file.filename.lower().endswith(".pdf"):
        return render_template("pdf_split.html", error=t("err.invalid_pdf"))

    if not range_str:
        return render_template("pdf_split.html", error=t("err.enter_page_range"))

    try:
        reader = PdfReader(file.stream)
        max_pages = len(reader.pages)
        
        # Parse range
        pages_to_keep = set()
        for part in range_str.split(','):
            part = part.strip()
            if '-' in part:
                start, end = part.split('-')
                start_idx = int(start.strip()) - 1
                end_idx = int(end.strip()) - 1
                for p in range(max(0, start_idx), min(max_pages, end_idx + 1)):
                    pages_to_keep.add(p)
            else:
                p_idx = int(part) - 1
                if 0 <= p_idx < max_pages:
                    pages_to_keep.add(p_idx)
                    
        sorted_pages = sorted(list(pages_to_keep))
        if not sorted_pages:
            return render_template("pdf_split.html", error=t("err.no_valid_pages_in_range"))

        writer = PdfWriter()
        for p in sorted_pages:
            writer.add_page(reader.pages[p])

        output = BytesIO()
        writer.write(output)
        output.seek(0)

        _log_action("pdf_split")
        return send_file(
            output,
            mimetype="application/pdf",
            as_attachment=True,
            download_name="toolboxquick-split.pdf",
        )
    except Exception as e:
        return render_template("pdf_split.html", error=t("err.generic_error", error=str(e)))


@app.route("/audio-tools/trim")
def audio_trim():
    return render_template("audio_trim.html")


@app.route("/audio-tools/record")
def audio_record():
    return render_template("audio_record.html")


@app.route("/video-tools/screen-record")
def video_screen_record():
    return render_template("video_screen_record.html")


@app.route("/convert/image")
def convert_image():
    return render_template("convert_image.html")


@app.route("/convert/image-compress")
def convert_image_compress():
    return render_template("image_compress.html")


# Temporary directory setup for ZIP extractor
TEMP_DIR = os.path.join(tempfile.gettempdir(), "toolboxquick_zips")
os.makedirs(TEMP_DIR, exist_ok=True)


@app.route("/pdf-tools/protect", methods=["GET", "POST"])
def pdf_protect():
    if request.method == "GET":
        return render_template("pdf_protect.html", error=None)

    file = request.files.get("pdf_file")
    password = request.form.get("password")

    if not file or not file.filename.lower().endswith(".pdf"):
        return render_template("pdf_protect.html", error=t("err.invalid_pdf"))

    if not password:
        return render_template("pdf_protect.html", error=t("err.set_password"))

    try:
        reader = PdfReader(file.stream)
        writer = PdfWriter()

        for page in reader.pages:
            writer.add_page(page)

        writer.encrypt(password)

        output = BytesIO()
        writer.write(output)
        output.seek(0)

        _log_action("pdf_protect")
        return send_file(
            output,
            mimetype="application/pdf",
            as_attachment=True,
            download_name="toolboxquick-protected.pdf"
        )
    except Exception as e:
        return render_template("pdf_protect.html", error=t("err.generic_error", error=str(e)))


@app.route("/pdf-tools/unlock", methods=["GET", "POST"])
def pdf_unlock():
    if request.method == "GET":
        return render_template("pdf_unlock.html", error=None)

    file = request.files.get("pdf_file")
    password = request.form.get("password")

    if not file or not file.filename.lower().endswith(".pdf"):
        return render_template("pdf_unlock.html", error=t("err.invalid_pdf"))

    if not password:
        return render_template("pdf_unlock.html", error=t("err.enter_password"))

    try:
        reader = PdfReader(file.stream)

        if not reader.is_encrypted:
            return render_template("pdf_unlock.html", error=t("err.pdf_not_encrypted"))

        decryption_result = reader.decrypt(password)
        if decryption_result == 0:
            return render_template("pdf_unlock.html", error=t("err.wrong_password"))

        writer = PdfWriter()
        for page in reader.pages:
            writer.add_page(page)

        output = BytesIO()
        writer.write(output)
        output.seek(0)

        _log_action("pdf_unlock")
        return send_file(
            output,
            mimetype="application/pdf",
            as_attachment=True,
            download_name="toolboxquick-unlocked.pdf"
        )
    except Exception as e:
        return render_template("pdf_unlock.html", error=t("err.generic_error", error=str(e)))


@app.route("/audio-tools/volume")
def audio_volume():
    return render_template("audio_volume.html")


@app.route("/audio-tools/speed")
def audio_speed():
    return render_template("audio_speed.html")


@app.route("/video-tools/rotate")
def video_rotate():
    return render_template("video_rotate.html")


@app.route("/video-tools/trim")
def video_trim():
    return render_template("video_trim.html")


@app.route("/video-tools/merge")
def video_merge():
    return render_template("video_merge.html")


@app.route("/convert/audio")
def convert_audio():
    return render_template("convert_audio.html")


@app.route("/convert/video")
def convert_video():
    return render_template("convert_video.html")


@app.route("/convert/extract", methods=["GET", "POST"])
def convert_extract():
    if request.method == "GET":
        return render_template("convert_extract.html", files=None, error=None)

    file = request.files.get("zip_file")
    if not file or not file.filename.lower().endswith(".zip"):
        return render_template("convert_extract.html", files=None, error=t("err.invalid_zip"))

    try:
        zip_id = str(uuid.uuid4())
        zip_path = os.path.join(TEMP_DIR, f"{zip_id}.zip")
        file.save(zip_path)

        files_list = []
        with zipfile.ZipFile(zip_path, 'r') as zip_ref:
            ok, err = _check_zip_safety(zip_ref)
            if not ok:
                os.remove(zip_path)
                return render_template("convert_extract.html", files=None, error=err)

            for info in zip_ref.infolist():
                if not info.is_dir():
                    files_list.append({
                        "name": info.filename,
                        "size": info.file_size
                    })

        _log_action("zip_extract")
        return render_template(
            "convert_extract.html",
            files=files_list,
            zip_id=zip_id,
            archive_name=file.filename,
            error=None
        )
    except Exception as e:
        return render_template("convert_extract.html", files=None, error=t("err.generic_error", error=str(e)))


@app.route("/convert/extract/download")
def convert_extract_download():
    zip_id = request.args.get("id")
    file_path = request.args.get("file")

    if not zip_id or not file_path:
        return t("err.invalid_request"), 400

    try:
        uuid.UUID(zip_id)
    except ValueError:
        return t("err.invalid_request"), 400

    zip_path = os.path.join(TEMP_DIR, f"{zip_id}.zip")
    if not os.path.exists(zip_path):
        return t("err.archive_not_found_or_expired"), 404

    try:
        with zipfile.ZipFile(zip_path, 'r') as zip_ref:
            if file_path not in zip_ref.namelist():
                return t("err.file_not_found_in_archive"), 404

            data = zip_ref.read(file_path)
            io_buf = BytesIO(data)
            basename = os.path.basename(file_path)
            
            return send_file(
                io_buf,
                as_attachment=True,
                download_name=basename
            )
    except Exception as e:
        return t("err.download_error", error=str(e)), 500


# --- AUDIO TOOLS ---
@app.route("/audio-tools/pitch")
def audio_pitch():
    return render_template("audio_pitch.html")

@app.route("/audio-tools/equalizer")
def audio_equalizer():
    return render_template("audio_equalizer.html")

@app.route("/audio-tools/joiner")
def audio_joiner():
    return render_template("audio_joiner.html")

@app.route("/audio-tools/reverse")
def audio_reverse():
    return render_template("audio_reverse.html")


# --- PDF TOOLS ---
@app.route("/pdf-tools/rotate")
def pdf_rotate_page():
    return render_template("pdf_rotate.html")

@app.route("/pdf-tools/page-numbers")
def pdf_page_numbers():
    return render_template("pdf_page_numbers.html")

@app.route("/pdf-tools/compress")
def pdf_compress_page():
    return render_template("pdf_compress.html")

@app.route("/pdf-tools/pdf-to-jpg")
def pdf_to_jpg():
    return render_template("pdf_to_img.html")

@app.route("/pdf-tools/pdf-to-png")
def pdf_to_png():
    return render_template("pdf_to_img.html")

@app.route("/pdf-tools/jpg-to-pdf")
def jpg_to_pdf():
    return render_template("img_to_pdf.html")

@app.route("/pdf-tools/png-to-pdf")
def png_to_pdf():
    return render_template("img_to_pdf.html")

@app.route("/pdf-tools/pdf-to-html")
def pdf_to_html():
    return render_template("pdf_to_html.html")

@app.route("/pdf-tools/word-to-pdf")
def word_to_pdf():
    return render_template("pdf_office.html", tool_type="word-pdf", tool_title=t("office.word_to_pdf.title"), tool_desc=t("office.word_to_pdf.desc"), file_accept=".docx")

@app.route("/pdf-tools/pdf-to-word", methods=["GET", "POST"])
def pdf_to_word():
    if request.method == "GET":
        return render_template("pdf_office.html", tool_type="pdf-word", tool_title=t("office.pdf_to_word.title"), tool_desc=t("office.pdf_to_word.desc"), file_accept=".pdf")

    file = request.files.get("pdf_file")
    if not file or not file.filename.lower().endswith(".pdf"):
        return jsonify({"error": t("err.invalid_pdf")}), 400

    try:
        output = BytesIO()
        cv = PdfToDocxConverter(stream=file.read())
        cv.convert(output)
        cv.close()
        output.seek(0)

        _log_action("pdf_to_word")
        return send_file(
            output,
            mimetype="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            as_attachment=True,
            download_name="toolboxquick-converted.docx",
        )
    except Exception as e:
        return jsonify({"error": t("err.generic_error", error=str(e))}), 500

@app.route("/pdf-tools/excel-to-pdf")
def excel_to_pdf():
    return render_template("pdf_office.html", tool_type="excel-pdf", tool_title=t("office.excel_to_pdf.title"), tool_desc=t("office.excel_to_pdf.desc"), file_accept=".xlsx,.xls")

@app.route("/pdf-tools/pdf-to-excel")
def pdf_to_excel():
    return render_template("pdf_office.html", tool_type="pdf-excel", tool_title=t("office.pdf_to_excel.title"), tool_desc=t("office.pdf_to_excel.desc"), file_accept=".pdf")

@app.route("/pdf-tools/ppt-to-pdf", methods=["GET", "POST"])
def ppt_to_pdf():
    if request.method == "GET":
        return render_template("pdf_office.html", tool_type="ppt-pdf", tool_title=t("office.ppt_to_pdf.title"), tool_desc=t("office.ppt_to_pdf.desc"), file_accept=".pptx")

    file = request.files.get("pptx_file")
    if not file or not file.filename.lower().endswith(".pptx"):
        return jsonify({"error": t("err.invalid_pptx")}), 400

    try:
        output = _convert_pptx_to_pdf(BytesIO(file.read()))

        _log_action("ppt_to_pdf")
        return send_file(
            output,
            mimetype="application/pdf",
            as_attachment=True,
            download_name="toolboxquick-converted.pdf",
        )
    except Exception as e:
        return jsonify({"error": t("err.generic_error", error=str(e))}), 500


# --- VIDEO TOOLS ---
@app.route("/video-tools/editor")
@app.route("/video-tools/crop")
@app.route("/video-tools/loop")
@app.route("/video-tools/video-volume")
@app.route("/video-tools/video-speed")
@app.route("/video-tools/add-text")
@app.route("/video-tools/add-image")
@app.route("/video-tools/stabilize")
@app.route("/video-tools/remove-logo")
@app.route("/video-tools/add-audio")
@app.route("/video-tools/resize")
def video_editor_tool():
    path_to_tab = {
        "video-volume": "volume",
        "video-speed": "speed",
        "add-text": "text",
        "add-image": "image",
    }
    slug = request.path.rsplit("/", 1)[-1]
    active_tool = path_to_tab.get(slug, slug)
    return render_template("video_editor.html", active_tool=active_tool)

@app.route("/video-tools/text-to-speech")
def video_text_to_speech():
    return render_template("video_text_to_speech.html")

@app.route("/video-tools/record-camera")
def video_record_camera():
    return render_template("video_record_camera.html")

@app.route("/video-tools/gif")
def video_gif():
    return render_template("video_gif.html")


# --- UTILITY TOOLS ---
@app.route("/qr-kod")
def qr_generator():
    return render_template("qr_generator.html")


# --- CONVERTERS ---
@app.route("/convert/document")
def convert_document():
    return render_template("convert_document.html")

@app.route("/convert/ebook")
def convert_ebook():
    return render_template("convert_ebook.html")

@app.route("/convert/font")
def convert_font():
    return render_template("convert_font.html")

@app.route("/convert/archive", methods=["GET", "POST"])
def convert_archive():
    if request.method == "GET":
        return render_template("convert_archive.html", error=None)
    
    file = request.files.get("archive_file")
    target_format = request.form.get("target_format")
    
    if not file or not file.filename:
        return render_template("convert_archive.html", error=t("err.invalid_archive"))
        
    try:
        in_buf = BytesIO(file.read())
        out_buf = BytesIO()
        
        files_dict = {}
        filename = file.filename.lower()
        
        if filename.endswith(".zip"):
            with zipfile.ZipFile(in_buf, 'r') as z:
                ok, err = _check_zip_safety(z)
                if not ok:
                    return render_template("convert_archive.html", error=err)
                for name in z.namelist():
                    files_dict[name] = z.read(name)
        elif filename.endswith(".tar") or filename.endswith(".gz") or filename.endswith(".tgz"):
            in_buf.seek(0)
            mode = "r:gz" if (filename.endswith(".gz") or filename.endswith(".tgz")) else "r:"
            with tarfile.open(fileobj=in_buf, mode=mode) as t:
                for member in t.getmembers():
                    if member.isfile():
                        f = t.extractfile(member)
                        if f:
                            files_dict[member.name] = f.read()
        else:
            return render_template("convert_archive.html", error=t("err.unsupported_archive_format"))
            
        out_buf.seek(0)
        if target_format == "zip":
            with zipfile.ZipFile(out_buf, 'w', zipfile.ZIP_DEFLATED) as z:
                for name, data in files_dict.items():
                    z.writestr(name, data)
            download_name = "converted.zip"
            mimetype = "application/zip"
        elif target_format == "tar":
            with tarfile.open(fileobj=out_buf, mode="w:") as t:
                for name, data in files_dict.items():
                    info = tarfile.TarInfo(name=name)
                    info.size = len(data)
                    t.addfile(tarinfo=info, fileobj=BytesIO(data))
            download_name = "converted.tar"
            mimetype = "application/x-tar"
        elif target_format == "tgz":
            with tarfile.open(fileobj=out_buf, mode="w:gz") as t:
                for name, data in files_dict.items():
                    info = tarfile.TarInfo(name=name)
                    info.size = len(data)
                    t.addfile(tarinfo=info, fileobj=BytesIO(data))
            download_name = "converted.tar.gz"
            mimetype = "application/gzip"
            
        out_buf.seek(0)
        _log_action(f"archive_convert_to_{target_format}")
        return send_file(
            out_buf,
            mimetype=mimetype,
            as_attachment=True,
            download_name=download_name
        )
    except Exception as e:
        return render_template("convert_archive.html", error=t("err.generic_error", error=str(e)))
